# meta developer: @B_Mods 
from .. import loader, utils
import pptx
import os

class PresText(loader.Module):
    """Присылает текст слайдов из PPTX"""
    
    strings = {"name": "PresText"}

    @loader.command()
    async def pres(self, message):
        """Отправляет текст со слайдов .pptx — использовать как ответ на файл"""
        
        reply = message.reply_to_message
        if not reply:
            return await message.edit("❗ Пришли .pptx файл и ответь командой `.pres`")

        # Проверяем файл
        if not reply.document:
            return await message.edit("❗ Это не файл. Пришли презентацию в формате .pptx")

        file = await reply.download()
        if not file.endswith(".pptx"):
            return await message.edit("❗ Нужен файл презентации .pptx")

        await message.edit("📥 Загружаю презентацию...")

        try:
            prs = pptx.Presentation(file)
        except Exception as e:
            return await message.edit(f"❌ Ошибка при чтении файла: {e}")

        await message.edit("📄 Извлекаю текст со слайдов...")

        if len(prs.slides) == 0:
            return await message.edit("❗ В презентации нет слайдов")

        # Обрабатываем каждый слайд
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        slide_text.append(txt)

            if not slide_text:
                text = "— Слайд без текста —"
            else:
                text = "\n".join(slide_text)

            await message.client.send_message(
                message.chat_id,
                f"📌 *Слайд {i}*\n\n{text}"
            )

        await message.edit("✅ Готово! Все тексты слайдов отправлены.")

        try:
            os.remove(file)
        except:
            pass
